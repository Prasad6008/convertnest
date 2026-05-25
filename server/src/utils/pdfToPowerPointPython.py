import sys
import os
import tempfile
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def color_int_to_rgb(color_value):
    if color_value is None:
        return RGBColor(0, 0, 0)

    r = (color_value >> 16) & 255
    g = (color_value >> 8) & 255
    b = color_value & 255

    return RGBColor(r, g, b)


def safe_inch(value):
    return Inches(max(0.01, value))


def exact_pdf_to_pptx(input_pdf, output_pptx):
    pdf = fitz.open(input_pdf)

    if len(pdf) == 0:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(output_pptx)
        return

    first_page = pdf[0]
    slide_width_in = first_page.rect.width / 72
    slide_height_in = first_page.rect.height / 72

    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as temp_dir:
        for index, page in enumerate(pdf):
            slide = prs.slides.add_slide(blank_layout)

            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = os.path.join(temp_dir, f"page-{index + 1}.png")
            pix.save(image_path)

            slide.shapes.add_picture(
                image_path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

    pdf.close()
    prs.save(output_pptx)


def editable_pdf_to_pptx(input_pdf, output_pptx):
    pdf = fitz.open(input_pdf)

    if len(pdf) == 0:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(output_pptx)
        return

    first_page = pdf[0]
    base_width_in = first_page.rect.width / 72
    base_height_in = first_page.rect.height / 72

    prs = Presentation()
    prs.slide_width = Inches(base_width_in)
    prs.slide_height = Inches(base_height_in)

    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as temp_dir:
        for page_index, page in enumerate(pdf):
            slide = prs.slides.add_slide(blank_layout)

            page_width = page.rect.width
            page_height = page.rect.height

            scale_x = base_width_in / (page_width / 72)
            scale_y = base_height_in / (page_height / 72)

            page_dict = page.get_text("dict")

            # Add image blocks first
            for block_index, block in enumerate(page_dict.get("blocks", [])):
                if block.get("type") == 1 and block.get("image"):
                    try:
                        x0, y0, x1, y1 = block["bbox"]
                        ext = block.get("ext", "png")
                        image_path = os.path.join(
                            temp_dir,
                            f"page-{page_index + 1}-image-{block_index}.{ext}"
                        )

                        with open(image_path, "wb") as image_file:
                            image_file.write(block["image"])

                        slide.shapes.add_picture(
                            image_path,
                            safe_inch((x0 / 72) * scale_x),
                            safe_inch((y0 / 72) * scale_y),
                            width=safe_inch(((x1 - x0) / 72) * scale_x),
                            height=safe_inch(((y1 - y0) / 72) * scale_y)
                        )
                    except Exception:
                        pass

            # Add editable text line by line
            for block in page_dict.get("blocks", []):
                if block.get("type") != 0:
                    continue

                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    line_text = "".join(span.get("text", "") for span in spans).strip()

                    if not line_text:
                        continue

                    x0, y0, x1, y1 = line["bbox"]

                    left = (x0 / 72) * scale_x
                    top = (y0 / 72) * scale_y
                    width = max(0.1, ((x1 - x0) / 72) * scale_x + 0.05)
                    height = max(0.1, ((y1 - y0) / 72) * scale_y + 0.04)

                    textbox = slide.shapes.add_textbox(
                        safe_inch(left),
                        safe_inch(top),
                        safe_inch(width),
                        safe_inch(height)
                    )

                    text_frame = textbox.text_frame
                    text_frame.clear()
                    text_frame.margin_left = 0
                    text_frame.margin_right = 0
                    text_frame.margin_top = 0
                    text_frame.margin_bottom = 0
                    text_frame.word_wrap = False

                    paragraph = text_frame.paragraphs[0]
                    paragraph.space_after = Pt(0)
                    paragraph.space_before = Pt(0)

                    for span in spans:
                        text = span.get("text", "")

                        if not text:
                            continue

                        run = paragraph.add_run()
                        run.text = text

                        font = run.font
                        font_name = span.get("font", "Arial")
                        font_size = span.get("size", 12)

                        font.name = font_name
                        font.size = Pt(max(4, min(72, font_size * scale_y)))
                        font.bold = "Bold" in font_name or "bold" in font_name
                        font.italic = "Italic" in font_name or "italic" in font_name
                        font.color.rgb = color_int_to_rgb(span.get("color"))

    pdf.close()
    prs.save(output_pptx)


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python pdfToPowerPointPython.py input.pdf output.pptx mode", file=sys.stderr)
        sys.exit(1)

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    mode = sys.argv[3]

    try:
        if mode == "editable":
            editable_pdf_to_pptx(input_pdf, output_pptx)
        else:
            exact_pdf_to_pptx(input_pdf, output_pptx)

        print("PDF to PowerPoint conversion completed")
    except Exception as error:
        print(str(error), file=sys.stderr)
        sys.exit(1)